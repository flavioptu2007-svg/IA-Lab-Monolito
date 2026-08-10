#!/usr/bin/env python3
"""
standardize_fonts.py
Padroniza fontes, tamanhos e capitalização de todos os títulos em documentos .docx e .pptx.

Padrão definido:
  - Fonte: Calibri
  - Tamanho título: 16pt (Negrito)
  - Tamanho subtítulo: 13pt (Negrito)
  - Tamanho corpo do texto: 11pt (Normal)
  - Capitalização: sentence case (só primeira letra maiúscula + nomes próprios)
"""

import re
import sys
from pathlib import Path

from docx import Document
from docx.shared import Pt
from docx.shared import RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Pt as PptxPt

# === CONFIGURAÇÃO DO PADRÃO ===
FONT_NAME = "Calibri"
TITLE_SIZE = 16
SUBTITLE_SIZE = 13
BODY_SIZE = 11
TITLE_COLOR_DOCX = DocxRGBColor(0x1A, 0x1A, 0x2E)  # Azul escuro profissional
TITLE_COLOR_PPTX = PptxRGBColor(0x1A, 0x1A, 0x2E)
BODY_COLOR_DOCX = DocxRGBColor(0x33, 0x33, 0x33)  # Cinza escuro
BODY_COLOR_PPTX = PptxRGBColor(0x33, 0x33, 0x33)

# Nomes próprios que devem permanecer com inicial maiúscula
PROPER_NOUNS = {
    # Pessoas
    "Lutero",
    "Calvino",
    "Zwingli",
    "Napoleão",
    "Darwin",
    "Lenin",
    "Trotski",
    "Jesus",
    "Cristo",
    "Moisés",
    "Maomé",
    "Coraci",
    # Países / Regiões / Gentílicos
    "Brasil",
    "Portugal",
    "Espanha",
    "Espanhola",
    "França",
    "Alemanha",
    "Inglaterra",
    "Itália",
    "Roma",
    "Grécia",
    "Egito",
    "México",
    "Peru",
    "Haiti",
    "Europa",
    "Ásia",
    "África",
    "América",
    "Américas",
    "Oriente Médio",
    # Cidades
    "Paracatu",
    "Brasília",
    "São Paulo",
    "Rio de Janeiro",
    "Lisboa",
    "Wittenberg",
    "Genebra",
    "Paris",
    "Londres",
    # Instituições
    "BNCC",
    "CRMG",
    "SAEB",
    "SIMAVE",
    "CAED",
    "INEP",
    "MEC",
    "UFJF",
    "USP",
    "EJA",
    "GRP",
    "SEE",
    "CACILDA",
    # Civilizações / Povos
    "Maias",
    "Astecas",
    "Incas",
    "Pr-Colombianos",
    "Pré-Colombianos",
    "Maia",
    "Asteca",
    "Inca",
    "Cristianismo",
    "Catolicismo",
    "Protestantismo",
    "Contrarreforma",
    "Igreja Católica",
    # Eventos históricos
    "Reforma Protestante",
    "Revolução Francesa",
    "Revolução Russa",
    "Independência do Haiti",
    "Guerra da Independência Espanhola",
    # Meses e dias (sempre minúsculos em português)
    # Tecnologias / Marcas
    "ChatGPT",
    "Gemini",
    "YouTube",
    "Google",
    "IA",
    # Termos educacionais
    "História",
    "Comum",
    "Bimestral",
    "Trimestral",
}

# Palavras que NUNCA devem ser maiúsculas (exceto início de frase)
ALWAYS_LOWER = {
    "de",
    "do",
    "da",
    "dos",
    "das",
    "e",
    "ou",
    "a",
    "o",
    "em",
    "no",
    "na",
    "nos",
    "nas",
    "um",
    "uma",
    "uns",
    "umas",
    "para",
    "por",
    "com",
    "sem",
    "que",
    "se",
    "ao",
    "aos",
    "à",
    "às",
    "este",
    "esta",
    "estes",
    "estas",
    "esse",
    "essa",
    "esses",
    "essas",
    "isto",
    "isso",
    "mas",
    "como",
    "não",
    "sim",
    "mais",
    "muito",
    "pouco",
    "já",
    "também",
    "só",
    "apenas",
}


def sentence_case(text):
    """
    Aplica capitalização estilo frase:
    - Primeira letra da frase em maiúscula
    - Meses, dias da semana, profissões em minúsculas
    - Nomes próprios com inicial maiúscula
    - Demais palavras em minúsculas
    """
    if not text or len(text) < 2:
        return text

    # Separa em tokens mantendo pontuação e espaços
    tokens = re.split(r"(\s+|[.,;:!?()\-—])", text)
    result = []
    at_start = True

    for token in tokens:
        if not token or token.isspace() or token in ".,;:!?()\\-—":
            result.append(token)
            continue

        # Verifica se é nome próprio
        is_proper = False
        for proper in PROPER_NOUNS:
            if token.lower() == proper.lower() or token == proper:
                token = proper
                is_proper = True
                break

        if not is_proper:
            # Início de frase = maiúscula
            if at_start:
                token = token[0].upper() + token[1:].lower() if len(token) > 1 else token.upper()
                at_start = False
            else:
                token = token.lower()

        # Reseta at_start após pontuação final
        if result and result[-1] in ".!?":
            at_start = True

        result.append(token)

    return "".join(result)


def is_likely_title(text, paragraph_idx):
    """Heurística: primeiros parágrafos ou texto curto = título."""
    text = text.strip()
    if not text or len(text) < 3:
        return None

    if paragraph_idx < 3 and len(text) < 80:
        return "title"

    if text.isupper() and len(text) < 60:
        return "title"

    return None


def standardize_docx(filepath, dry_run=True):
    """Padroniza fontes e capitalização em documento .docx"""
    doc = Document(str(filepath))
    changed = 0
    para_idx = 0

    for para in doc.paragraphs:
        if not para.text.strip():
            continue

        text = para.text.strip()
        role = is_likely_title(text, para_idx)

        style_name = para.style.name.lower() if para.style else ""
        if "heading" in style_name or "título" in style_name or "title" in style_name:
            role = "title"
        elif "subtitle" in style_name:
            role = "subtitle"

        para_idx += 1

        if role is None:
            role = "body"

        # Define tamanho, cor e capitalização baseado no papel
        if role == "title":
            size = TITLE_SIZE
            bold = True
            color = TITLE_COLOR_DOCX
            new_text = sentence_case(text)
        elif role == "subtitle":
            size = SUBTITLE_SIZE
            bold = True
            color = TITLE_COLOR_DOCX
            new_text = sentence_case(text)
        else:
            size = BODY_SIZE
            bold = False
            color = BODY_COLOR_DOCX
            new_text = text

        # Verifica se o texto mudou
        text_changed = new_text != text

        # Aplica a cada run do parágrafo
        for run in para.runs:
            old_font = run.font.name
            old_size = run.font.size
            old_bold = run.bold
            old_color = run.font.color.rgb if run.font.color and run.font.color.rgb else None
            old_text = run.text

            run.font.name = FONT_NAME
            run.font.size = Pt(size)
            run.bold = bold
            run.font.color.rgb = color

            # Substitui texto se mudou capitalização (apenas no primeiro run)
            if text_changed and run == para.runs[0]:
                run.text = new_text

            if (
                old_font != FONT_NAME
                or old_size != Pt(size)
                or old_bold != bold
                or (old_color and str(old_color) != str(color))
                or (text_changed and old_text != new_text)
            ):
                changed += 1

    if not dry_run and changed > 0:
        doc.save(str(filepath))

    return changed


def standardize_pptx(filepath, dry_run=True):
    """Padroniza fontes e capitalização em apresentação .pptx"""
    prs = Presentation(str(filepath))
    changed = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                text = paragraph.text.strip()
                if not text:
                    continue

                role = "body"
                if para_idx == 0 and len(text) < 80:
                    role = "title"
                elif para_idx == 1 and len(text) < 60:
                    role = "subtitle"

                if shape.is_placeholder:
                    ph_type = str(shape.placeholder_format.type).lower()
                    if "title" in ph_type:
                        role = "title"
                    elif "subtitle" in ph_type:
                        role = "subtitle"

                if role == "title":
                    size = TITLE_SIZE
                    bold = True
                    color = TITLE_COLOR_PPTX
                    new_text = sentence_case(text)
                elif role == "subtitle":
                    size = SUBTITLE_SIZE
                    bold = True
                    color = TITLE_COLOR_PPTX
                    new_text = sentence_case(text)
                else:
                    size = BODY_SIZE
                    bold = False
                    color = BODY_COLOR_PPTX
                    new_text = text

                text_changed = new_text != text

                for run in paragraph.runs:
                    old_font = run.font.name
                    old_size = run.font.size
                    old_text = run.text

                    run.font.name = FONT_NAME
                    run.font.size = PptxPt(size)
                    run.font.bold = bold
                    run.font.color.rgb = color

                    if text_changed and run == paragraph.runs[0]:
                        run.text = new_text

                    if (
                        old_font != FONT_NAME
                        or old_size != PptxPt(size)
                        or (text_changed and old_text != new_text)
                    ):
                        changed += 1

    if not dry_run and changed > 0:
        prs.save(str(filepath))

    return changed


def process_all(dry_run=True):
    """Processa todos os documentos na estrutura Escola"""
    base = Path("/home/flavio/Documentos/Escola")
    docx_files = list(base.rglob("*.docx"))
    pptx_files = list(base.rglob("*.pptx"))
    odt_files = list(base.rglob("*.odt"))

    total_files = len(docx_files) + len(pptx_files)
    total_changes = 0

    print(f"Padronização de fontes: {FONT_NAME}")
    print(f"  Títulos:   {TITLE_SIZE}pt, Negrito, Cor escura")
    print(f"  Subtítulos: {SUBTITLE_SIZE}pt, Negrito, Cor escura")
    print(f"  Corpo:     {BODY_SIZE}pt, Normal, Cinza escuro")
    print(f"\nArquivos encontrados: {total_files} ({len(docx_files)} DOCX, {len(pptx_files)} PPTX)")
    if odt_files:
        print(f"  ODT: {len(odt_files)} (não suportado — use LibreOffice CLI)")
    print()

    for filepath in sorted(docx_files):
        try:
            changes = standardize_docx(filepath, dry_run=dry_run)
            if changes > 0:
                action = "PADRONIZAR" if not dry_run else "[DRY] PADRONIZAR"
                rel = filepath.relative_to(base)
                print(f"  {action}: {rel} ({changes} alterações)")
                total_changes += changes
        except Exception as e:
            print(f"  ERRO: {filepath.relative_to(base)} — {e}")

    for filepath in sorted(pptx_files):
        try:
            changes = standardize_pptx(filepath, dry_run=dry_run)
            if changes > 0:
                action = "PADRONIZAR" if not dry_run else "[DRY] PADRONIZAR"
                rel = filepath.relative_to(base)
                print(f"  {action}: {rel} ({changes} alterações)")
                total_changes += changes
        except Exception as e:
            print(f"  ERRO: {filepath.relative_to(base)} — {e}")

    print(f"\n{'=' * 60}")
    print(f"{'[DRY RUN — nada foi alterado]' if dry_run else '[EXECUTADO]'}")
    print(f"{'=' * 60}")
    print(f"  Total de alterações de fonte: {total_changes}")
    print(f"{'=' * 60}")
    if dry_run:
        print("\nRevise acima. Se estiver correto, rode com --apply")


if __name__ == "__main__":
    dry = "--apply" not in sys.argv
    process_all(dry_run=dry)
